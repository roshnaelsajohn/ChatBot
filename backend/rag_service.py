"""
RAG Service Module
Handles document embedding, storage in ChromaDB, and retrieval.
"""

# Fix for SSL Certificate Verify Failed error
# We rely on environment variables (HF_HUB_DISABLE_SSL_VERIFY) to handle this.
# This simple check ensures we respect the disable flag if set.
import os

if os.environ.get("HF_HUB_DISABLE_SSL_VERIFY") == "1":
    import ssl
    try:
        _create_unverified_https_context = ssl._create_unverified_context
        ssl._create_default_https_context = _create_unverified_https_context
    except AttributeError:
        pass

import hashlib
import re
from sentence_transformers import SentenceTransformer
import chromadb
from langchain_text_splitters import RecursiveCharacterTextSplitter, MarkdownHeaderTextSplitter, HTMLHeaderTextSplitter

# Adapter to use SentenceTransformer with LangChain components
class EmbeddingsAdapter:
    def __init__(self, model):
        self.model = model
    def embed_documents(self, texts):
        return self.model.encode(texts, normalize_embeddings=True).tolist()
    def embed_query(self, text):
        return self.model.encode([text], normalize_embeddings=True)[0].tolist()

class RAGService:
    """Service class for RAG operations with ChromaDB and Nomic embeddings."""
    
    def __init__(self, persist_directory: str = "./chroma_db"):
        # ... (init code remains overlapping with original, only showing changes implicitly via replacement context if needed, but here I am inserting methods)
        self.persist_directory = persist_directory
        self.client = chromadb.PersistentClient(path=persist_directory)
        self.embedding_model = SentenceTransformer("nomic-ai/nomic-embed-text-v1.5", trust_remote_code=True)
        self.embeddings_adapter = EmbeddingsAdapter(self.embedding_model)
        
        self.collection_name = "documents_nomic"
        try:
            self.collection = self.client.get_or_create_collection(name=self.collection_name, metadata={"hnsw:space": "cosine"})
        except Exception as e:
            print(f"Warning initializing collection: {e}")
            self.collection = self.client.get_or_create_collection(name=self.collection_name, metadata={"hnsw:space": "cosine"})
        
        # Base splitter for fallbacks
        self.text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )

    def split_document(self, text: str, file_type: str) -> list[str]:
        """
        Split document using advanced strategies based on file type.
        """
        chunks = []
        try:
            # 1. Structured Files (MD) -> Section Aware
            if file_type == 'md':
                headers_to_split_on = [("#", "Header 1"), ("##", "Header 2"), ("###", "Header 3")]
                splitter = MarkdownHeaderTextSplitter(headers_to_split_on=headers_to_split_on)
                # Split based on headers
                docs = splitter.split_text(text)
                
                # Further split large sections using recursive splitter to ensure size limits
                # We need to treat 'docs' (Document objects) back to text or use split_documents
                final_docs = self.text_splitter.split_documents(docs)
                chunks = [doc.page_content for doc in final_docs]

            # 2. HTML -> Section Aware (if we had raw HTML, but here we likely have text extracted. 
            # If we want true HTML splitting we need the raw content before extraction. 
            # For now, if we have cleaned text, we might treat it as unstructured or use semantic.)
            # Assuming 'text' passed here is already cleaned/extracted text.
            # If we want true HTML splitting, we should change how extract_text works, but for this step:
            
            # 3. Unstructured Files -> Semantic Chunking
            elif file_type in ['pdf', 'txt', 'docx', 'pptx']:
                try:
                    from langchain_experimental.text_splitter import SemanticChunker
                    # Use percentile threshold 
                    semantic_splitter = SemanticChunker(
                        self.embeddings_adapter, 
                        breakpoint_threshold_type="percentile"
                    )
                    chunks = semantic_splitter.split_text(text)
                    
                    # Safety check: If semantic chunker returns huge chunks, sub-split them
                    # or if it returns nothing
                    if not chunks:
                         chunks = self.text_splitter.split_text(text)
                except ImportError:
                    print("langchain_experimental not found, falling back to recursive splitter.")
                    chunks = self.text_splitter.split_text(text)
                except Exception as e:
                    print(f"Semantic chunking failed: {e}. Falling back.")
                    chunks = self.text_splitter.split_text(text)
            
            else:
                chunks = self.text_splitter.split_text(text)
                
        except Exception as e:
            print(f"Error in split_document: {e}. Using fallback.")
            chunks = self.text_splitter.split_text(text)
            
        return chunks

    def clean_text(self, text: str) -> str:
        """Clean and normalize text."""
        # Use regex to replace multiple spaces/tabs/newlines with single space
        # But for tables, we might want to preserve structure.
        # Let's clean multiple spaces but preserve single newlines?
        # For now, standard RAG cleaning:
        text = re.sub(r'\s+', ' ', text).strip()
        return text

    def _convert_table_to_text(self, table: list) -> str:
        """Convert a list of lists (table) to a text representation (Header: Value)."""
        if not table:
            return ""
        
        text_output = ""
        # Filter out None/Empty
        cleaned_table = [[str(cell or "").replace("\n", " ").strip() for cell in row] for row in table]
        
        if not cleaned_table:
            return ""

        # Assume first row is header
        headers = cleaned_table[0]
        
        # Process Rows
        for i, row in enumerate(cleaned_table[1:]):
             row_text = []
             for j, cell in enumerate(row):
                 # Handle cases where row length > header length
                 header_name = headers[j] if j < len(headers) else f"Column {j+1}"
                 row_text.append(f"{header_name}: {cell}")
             
             text_output += f"Row {i+1}: " + ", ".join(row_text) + "\n"
             
        return text_output

    def extract_text_from_file(self, file_path: str, file_type: str) -> str:
        """Extract text content from various file types, handling tables."""
        text = ""
        try:
            if file_type == 'pdf':
                try:
                    import pdfplumber
                    with pdfplumber.open(file_path) as pdf:
                        for page in pdf.pages:
                            # 1. Extract Tables as Markdown (Preserves Structure)
                            tables = page.extract_tables()
                            if tables:
                                for table in tables:
                                    md_table = self._convert_table_to_text(table)
                                    text += f"\n[Table Data]:\n{md_table}\n"
                            
                            # 2. Extract Generic Text (Fallback/Context)
                            # Using layout=True helps preserve visual structure of non-table elements too
                            text += page.extract_text(layout=True) or "" + "\n"
                            
                except ImportError:
                    print("pdfplumber not found, falling back to pypdf")
                    from pypdf import PdfReader
                    reader = PdfReader(file_path)
                    for page in reader.pages:
                        text += page.extract_text() + "\n"
                        
            elif file_type == 'docx':
                 import docx
                 doc = docx.Document(file_path)
                 
                 # Extract paragraphs
                 for para in doc.paragraphs:
                     text += para.text + "\n"
                     
                 # Extract tables
                 # TODO: Interleave properly if possible, but appending is better than missing.
                 if doc.tables:
                     text += "\n[Extracted Tables]:\n"
                     for table in doc.tables:
                        # Convert docx table to list of lists
                        data = []
                        for row in table.rows:
                            data.append([cell.text for cell in row.cells])
                        text += self._convert_table_to_text(data) + "\n"

            elif file_type == 'pptx':
                from pptx import Presentation
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                        if shape.has_table:
                             table = shape.table
                             data = []
                             for row in table.rows:
                                 data.append([cell.text_frame.text for cell in row.cells])
                             text += "\n[Table Data]:\n" + self._convert_table_to_text(data) + "\n"

            elif file_type in ['txt', 'md', 'html']:
                # For HTML, usage of BS4 is recommended if available
                if file_type == 'html':
                    from bs4 import BeautifulSoup
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        soup = BeautifulSoup(f, 'html.parser')
                        text = soup.get_text(separator='\n')
                else:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text = f.read()

        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return ""
            
        return self.clean_text(text)

    def generate_embeddings(self, texts: list[str], is_document: bool = True) -> list:
        """Generate embeddings using Nomic Model with correct prefixes."""
        prefix = "search_document: " if is_document else "search_query: "
        # Nomic v1.5 requires prefixes
        prefixed_texts = [prefix + text for text in texts]
        return self.embedding_model.encode(prefixed_texts, normalize_embeddings=True).tolist()

    def process_and_store_document(self, file_path: str, filename: str, file_type: str) -> dict:
        """Process a document and store it in ChromaDB."""
        try:
            # Extract text
            text = self.extract_text_from_file(file_path, file_type)
            if not text.strip():
                return {"success": False, "message": "No text content found"}
                
            # Smart Split
            chunks = self.split_document(text, file_type)
            
            # Generate embeddings manually
            embeddings = self.generate_embeddings(chunks, is_document=True)
            
            # Prepare IDs and Metadata
            doc_id = hashlib.md5(filename.encode()).hexdigest()[:8]
            ids = [f"{doc_id}_chunk_{i}" for i in range(len(chunks))]
            metadatas = [
                {"source": filename, "chunk_index": i, "total_chunks": len(chunks)}
                for i in range(len(chunks))
            ]
            
            # Add to ChromaDB
            try:
                self.collection.add(
                    documents=chunks,
                    embeddings=embeddings,
                    ids=ids,
                    metadatas=metadatas
                )
            except Exception as e:
                # Handle dimension mismatch by resetting collection
                if "dimension" in str(e).lower():
                    print("Dimension mismatch detected. Resetting collection for Nomic V2.")
                    self.client.delete_collection(self.collection_name)
                    self.collection = self.client.create_collection(
                        name=self.collection_name,
                        metadata={"hnsw:space": "cosine"}
                    )
                    # Retry add
                    self.collection.add(
                        documents=chunks,
                        embeddings=embeddings,
                        ids=ids,
                        metadatas=metadatas
                    )
                else:
                    raise e
            
            return {
                "success": True, 
                "message": f"Successfully processed {filename} with Nomic Embeddings", 
                "chunks_added": len(chunks)
            }
            
        except Exception as e:
            return {"success": False, "message": f"Error: {str(e)}"}
    
    def rerank_results(self, results: list, query: str, threshold: float = 0.55) -> list:
        """
        Rerank and filter results based on Hybrid Search logic.
        - Boosts exact keyword matches for short queries.
        - Applies strict semantic threshold.
        """
        refined_results = []
        is_keyword_search = len(query.split()) < 3
        
        for r in results:
            score = round(r.get('similarity', 0.0), 3)
            content = r.get('content', '')
            has_keyword = query.lower() in content.lower()
            
            # DECISION LOGIC:
            if is_keyword_search:
                if has_keyword:
                    # Boost score for exact match (lower floor)
                    if score >= 0.35: 
                        refined_results.append(r)
                elif score >= 0.65:
                    # Keep high-confidence semantic matches even without keyword
                    refined_results.append(r)
            else:
                # Standard Semantic Search
                if score >= threshold:
                    refined_results.append(r)
                    
        return refined_results

    def query_documents(self, query: str, n_results: int = 3) -> dict:
        """Query the vector database."""
        try:
            # Generate embedding for query
            query_embedding = self.generate_embeddings([query], is_document=False)
            
            results = self.collection.query(
                query_embeddings=query_embedding,
                n_results=n_results
            )
            
            formatted_results = []
            if results and results["documents"]:
                for i, doc in enumerate(results["documents"][0]):
                    formatted_results.append({
                        "content": doc,
                        "source": results["metadatas"][0][i]["source"],
                        "distance": results["distances"][0][i] if results["distances"] else 1.0,
                        "similarity": 1.0 - (results["distances"][0][i] if results["distances"] else 1.0)
                    })
            
            return {"success": True, "results": formatted_results}
            
        except Exception as e:
            return {"success": False, "message": f"Error querying: {str(e)}"}
            
    def get_collection_stats(self) -> dict:
        """Get collection stats."""
        try:
            return {"total_chunks": self.collection.count()}
        except:
            return {"total_chunks": 0}

    def get_all_documents(self) -> list:
        """Get a list of all unique documents in the collection."""
        try:
            # We can't query distinctly easily with Chroma's basic API efficiently for large datasets,
            # but for a small app, we can fetch all metadatas.
            # A more efficient way for large datasets would be to maintain a separate 'files' collection or SQLite table.
            # For this MVP, we'll fetch all metadatas (limit to 10000 chunks for safety).
            result = self.collection.get(include=["metadatas"], limit=10000)
            seen_files = set()
            file_list = []
            
            if result and result["metadatas"]:
                for meta in result["metadatas"]:
                    source = meta.get("source")
                    if source and source not in seen_files:
                        seen_files.add(source)
                        file_list.append(source)
            
            return sorted(list(file_list))
        except Exception as e:
            print(f"Error fetching documents: {e}")
            return []    

    def document_exists(self, filename: str) -> bool:
        """Check if a document with the given filename already exists."""
        try:
            # Verify if any chunk has this source
            result = self.collection.get(
                where={"source": filename},
                limit=1
            )
            return len(result["ids"]) > 0
        except:
            return False

    def delete_document(self, filename: str) -> dict:
        """Delete a document by filename."""
        try:
            # Delete entries where source equals filename
            self.collection.delete(
                where={"source": filename}
            )
            return {"success": True, "message": f"Document '{filename}' deleted"}
        except Exception as e:
            return {"success": False, "message": f"Error deleting document: {str(e)}"}

    def clear_collection(self) -> dict:
        """Clear the collection."""
        try:
            self.client.delete_collection(self.collection_name)
            self.collection = self.client.create_collection(
                name=self.collection_name,
                metadata={"hnsw:space": "cosine"}
            )
            return {"success": True, "message": "Collection cleared"}
        except Exception as e:
            return {"success": False, "message": str(e)}


# Singleton instance
_rag_service = None

def get_rag_service() -> RAGService:
    """Get or create the RAG service singleton."""
    global _rag_service
    if _rag_service is None:
        _rag_service = RAGService()
    return _rag_service
